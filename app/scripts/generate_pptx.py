# generate_pptx_ultra_pro.py – FINAL VERSION

import os
import json
import pandas as pd
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# --- Paths ---
ANSWERS_CSV = "data/auto_answered_questions.csv"
CLASSIFIED_JSON = "data/classified_questions.json"
LOGO_PATH = "assets/shield.png"
OUTPUT_DIR = "output"
BASE_NAME = "due_diligence_report"

# --- Dynamic Output ---
def get_next_output_filename():
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    i = 0
    while True:
        suffix = f" ({i})" if i else ""
        path = os.path.join(OUTPUT_DIR, f"{BASE_NAME}{suffix}.pptx")
        if not os.path.exists(path):
            return path
        i += 1

# --- Objectives ---
TAG_OBJECTIVES = {
    "AML / KYC": "Ensure the company has robust Anti-Money Laundering (AML) and Know Your Customer (KYC) policies and compliance in place.",
    "Community & UX": "Assess user engagement, platform usability, and transparency of communication with the community.",
    "Custody & Asset Security": "Evaluate how the company stores and secures digital assets and whether custody mechanisms are robust.",
    "Cybersecurity & Data Privacy": "Examine the strength of cybersecurity practices and safeguards for user and company data.",
    "ESG & Sustainability": "Evaluate the company’s environmental, social, and governance practices and long-term sustainability impact.",
    "Financial Health": "Ensure the company’s financial statements, cash flows, and projections reflect a healthy financial position.",
    "Governance": "Assess decision-making structure, roles of stakeholders, and board oversight processes.",
    "IP & Contracts": "Evaluate the clarity of intellectual property (IP) ownership and the contractual integrity of the business.",
    "Legal & Regulatory": "Ensure the company’s legal entity, licensing, and jurisdictional compliance are clearly established.",
    "Risk Management": "Identify how the company manages internal and external risks, including business continuity.",
    "Strategy & Competitive Positioning": "Understand the company’s business model, growth strategy, and differentiation in the market.",
    "Technology & Infrastructure": "Evaluate the scalability, reliability, and modernity of the company’s technological backbone.",
    "Tokenomics & Trading Integrity": "Assess token issuance mechanisms, market manipulation controls, and valuation transparency."
}

# --- Data Loader ---
def load_data():
    df = pd.read_csv(ANSWERS_CSV)
    df.columns = [c.lower() for c in df.columns]
    df["question"] = df["question"].astype(str).str.strip()
    with open(CLASSIFIED_JSON, "r", encoding="utf-8") as f:
        classified = json.load(f)
    return df, classified

# --- Extract Key Point ---
def extract_summary(text):
    if not isinstance(text, str) or len(text.strip()) < 30:
        return None
    text = text.replace("Based on the document,", "").replace("According to", "").strip()
    return text.split(". ")[0].strip().capitalize() + "."

# --- Format Issues ---
def format_issue(text):
    text = str(text).strip().rstrip(".")
    if not text:
        return None
    keywords = ["AML", "regulatory", "audit", "compliance", "governance"]
    if any(k.lower() in text.lower() for k in keywords):
        return f"Potential weakness: {text}"
    return f"Insufficient detail on {text.lower()}"

# --- Branding ---
def add_branding(slide):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(0.3), Inches(0.1), height=Inches(0.5))
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(8), Inches(0.3)).text_frame
    footer.text = "DueXpert – AI Crypto Fund Due Diligence Suite | info@duexpert.ai | www.duexpert.ai"
    para = footer.paragraphs[0]
    para.font.size = Pt(9)
    para.font.color.rgb = RGBColor(130, 130, 130)
    para.font.name = "Segoe UI"

# --- Slide Builder ---
def add_slide(prs, tag, findings, gaps, page_number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = 0.8

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8), Inches(0.5)).text_frame
    title_box.text = f"{page_number}. {tag}"
    p = title_box.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.name = "Segoe UI"

    # Objective
    y += 0.7
    obj_text = TAG_OBJECTIVES.get(tag, "[Objective not defined]")
    obj_box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8), Inches(0.5)).text_frame
    obj_box.text = f"Objective: {obj_text}"
    para = obj_box.paragraphs[0]
    para.font.size = Pt(13)
    para.font.italic = True
    para.font.name = "Segoe UI"
    para.font.color.rgb = RGBColor(90, 90, 90)

    # Key Findings
    y += 0.8
    findings_box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.5), Inches(2)).text_frame
    findings_box.text = "Key Findings:"
    findings_box.paragraphs[0].font.size = Pt(14)
    findings_box.paragraphs[0].font.bold = True
    findings_box.paragraphs[0].font.name = "Segoe UI"

    for f in findings[:5]:
        para = findings_box.add_paragraph()
        para.text = f"• {f}"
        para.font.size = Pt(12)
        para.font.name = "Segoe UI"

    # Outstanding Issues
    y += 2.3
    issues_box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.5), Inches(2)).text_frame
    issues_box.text = "Outstanding Issues:"
    issues_box.paragraphs[0].font.size = Pt(14)
    issues_box.paragraphs[0].font.bold = True
    issues_box.paragraphs[0].font.name = "Segoe UI"

    for g in gaps[:5]:
        para = issues_box.add_paragraph()
        para.text = f"• {format_issue(g)}"
        para.font.size = Pt(12)
        para.font.name = "Segoe UI"

    # Page number
    pg = slide.shapes.add_textbox(Inches(8.2), Inches(6.8), Inches(1), Inches(0.3)).text_frame
    pg.text = f"Page {page_number}"
    pg.paragraphs[0].font.size = Pt(10)
    pg.paragraphs[0].font.name = "Segoe UI"
    pg.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)

    add_branding(slide)

# --- Main Controller ---
def main():
    output_path = get_next_output_filename()
    prs = Presentation()
    df, questions = load_data()

    # --- Cover ---
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = "Due Diligence Report"
    cover.placeholders[1].text = f"Prepared by DueXpert | {datetime.today().strftime('%Y-%m-%d')}"
    add_branding(cover)

    # --- TOC ---
    toc = prs.slides.add_slide(prs.slide_layouts[6])
    toc_title = toc.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8), Inches(0.5)).text_frame
    toc_title.text = "CONTENTS"
    toc_title.paragraphs[0].font.size = Pt(24)
    toc_title.paragraphs[0].font.bold = True
    toc_title.paragraphs[0].font.name = "Segoe UI"

    from collections import defaultdict
    tag_map = defaultdict(list)
    for q in questions:
        tag_map[q["tag"]].append(q)

    toc_entries = []
    page_counter = 3
    for idx, tag in enumerate(sorted(tag_map), 1):
        toc_entries.append((idx, tag.upper(), page_counter))
        page_counter += 1

    for i, (idx, label, page) in enumerate(toc_entries):
        y = 1.0 + i * 0.4
        left = toc.shapes.add_textbox(Inches(0.7), Inches(y), Inches(6), Inches(0.4)).text_frame
        left.text = f"{idx}. {label}"
        left.paragraphs[0].font.size = Pt(12)
        left.paragraphs[0].font.name = "Segoe UI"

        right = toc.shapes.add_textbox(Inches(7.5), Inches(y), Inches(1), Inches(0.4)).text_frame
        right.text = str(page)
        right.paragraphs[0].font.size = Pt(12)
        right.paragraphs[0].font.name = "Segoe UI"

    # --- Per Domain Slide ---
    page_number = 3
    for tag in sorted(tag_map):
        findings, gaps = [], []
        for q in tag_map[tag]:
            qt = q["question"]
            row = df[df["question"] == qt]
            if not row.empty:
                summary = extract_summary(row.iloc[0]["answer"])
                if summary:
                    findings.append(summary)
                else:
                    gaps.append(qt)
            else:
                gaps.append(qt)
        add_slide(prs, tag, findings, gaps, page_number)
        page_number += 1

    prs.save(output_path)
    print(f"✅ FINAL ULTRA PRO REPORT READY: {output_path}")

if __name__ == "__main__":
    main()
