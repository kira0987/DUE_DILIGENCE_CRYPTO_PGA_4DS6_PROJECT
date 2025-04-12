from datetime import datetime  
from pptx import Presentation
from pptx.util import Inches, Pt
import io
import plotly.graph_objects as go

def generate_ppt(qa_results: list, risk_score: float, risk_terms: dict, entities: dict) -> bytes:
    prs = Presentation()
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Crypto Fund Due Diligence Report"
    subtitle.text = f"Risk Score: {risk_score}/100 | Generated on {datetime.now().date()}"  

    # --- Risk Overview Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Risk Analysis"
    
    # Add risk score visualization
    fig = go.Figure(go.Indicator(
        mode="gauge+number",
        value=risk_score,
        domain={'x': [0, 1], 'y': [0, 1]},
        gauge={'axis': {'range': [0, 100]},
               'steps': [
                   {'range': [0, 40], 'color': "green"},
                   {'range': [40, 70], 'color': "orange"},
                   {'range': [70, 100], 'color': "red"}]}))
    
    img_stream = io.BytesIO()
    fig.write_image(img_stream, format='png')
    slide.shapes.add_picture(img_stream, Inches(1), Inches(2), width=Inches(4))

    # --- Q&A Summary Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Q&A Insights"
    
    content = ""
    for i, qa in enumerate(qa_results[:5]):  # Top 5 Q&A
        content += f"{i+1}. {qa['question']}\n- {qa['answer'][:100]}...\n\n"
    
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), width=Inches(9), height=Inches(4))
    tf = textbox.text_frame
    tf.text = content

    # Save to bytes for Streamlit download
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    return ppt_stream.getvalue()