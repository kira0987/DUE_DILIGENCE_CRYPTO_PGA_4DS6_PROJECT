import os
import numpy as np
from pptx import Presentation
import io
os.environ["STREAMLIT_SERVER_ENABLE_FILE_WATCHER"] = "false"
os.environ['KMP_DUPLICATE_LIB_OK'] = 'TRUE'
np.float = np.float32
np.int = np.int32
import json
import uuid
import streamlit as st
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime
from pathlib import Path
import pandas as pd
from modules.pdf_text_extractor import extract_text_from_pdf
from modules.text_cleaning import clean_text
from modules.text_chunker import SmartChunker
from modules.vector_database import VectorDatabase
from modules.question_bank_loader import convert_pdf_to_question_bank
from modules.feature_extraction import extract_risk_features, extract_named_entities
from modules.llm_wrapper import llama_chat, generate_concise_prompt, evaluate_answer
from modules.data_loader import DataLoader
from modules.market_analyzer import MarketAnalyzer
from modules.pptx_report import generate_ppt

# Initialize components
vector_db = VectorDatabase()
data_loader = DataLoader()
market_analyzer = MarketAnalyzer()

# ---------------------------
# 🔧 Streamlit Setup
# ---------------------------
st.set_page_config(page_title="DueXpert", layout="wide")
st.title("📄 DueXpert: Crypto Fund Due Diligence")

st.sidebar.header("⚙️ Settings")
MAX_QUESTIONS = st.sidebar.slider(
    "How many questions to auto-answer?", min_value=1, max_value=100, value=5
)

UPLOAD_DIR = "user_uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

# ---------------------------
# 📊 Market Data Loading and Analysis
# ---------------------------
@st.cache_data
def load_and_analyze_market_data():
    """Load and analyze market data (cached for performance)"""
    market_data = data_loader.load_market_data()
    return market_analyzer.analyze_all_markets(market_data)

# Load market data at startup
market_analysis = load_and_analyze_market_data()

def display_market_insights():
    """Display key market insights in the UI"""
    st.sidebar.subheader("📈 Market Overview")
    
    current_entities = st.session_state.get('current_entities', {})
    crypto_projects = current_entities.get('crypto_project', [])
    
    if crypto_projects:
        display_assets = []
        for project in crypto_projects:
            for asset in market_analysis:
                if any(word.lower() in asset.lower() for word in project.split()):
                    display_assets.append(asset)
        display_assets = list(set(display_assets))[:3]
    else:
        display_assets = sorted(
            market_analysis.items(),
            key=lambda x: x[1].get('avg_daily_volume', 0),
            reverse=True
        )[:3]
        display_assets = [asset[0] for asset in display_assets]
    
    for asset in display_assets:
        metrics = market_analysis.get(asset, {})
        if metrics:
            col1, col2 = st.sidebar.columns([1, 2])
            with col1:
                st.metric(f"{asset} Price", f"${metrics.get('current_price', 0):,.2f}")
            with col2:
                st.metric("30D Vol", f"{metrics.get('30_day_volatility', 0):.2%}")

# ---------------------------
# 📚 Load Question Bank
# ---------------------------
QUESTION_BANK_PDF = "data/question_bank.pdf"
QUESTION_BANK_JSON = "data/question_bank.json"
if not os.path.exists(QUESTION_BANK_JSON):
    convert_pdf_to_question_bank(QUESTION_BANK_PDF, QUESTION_BANK_JSON)

with open(QUESTION_BANK_JSON, "r", encoding="utf-8") as f:
    question_bank = json.load(f)

# ---------------------------
# 🤖 Enhanced Q&A Functions with Market Context
# ---------------------------
def process_question_batch(questions: list, index_path: str, doc_metadata: dict) -> list:
    """Process questions in optimized batches with enhanced context"""
    if not vector_db.is_trusted_path(index_path):
        return [{
            "question": q["question"],
            "answer": "⚠️ Security violation detected",
            "category": q.get("category", "General")
        } for q in questions]
    
    results = []
    with ThreadPoolExecutor(max_workers=5) as executor:
        futures = []
        for q in questions:
            futures.append(executor.submit(
                process_single_question,
                q,
                index_path,
                {
                    **doc_metadata,
                    "market_analysis": market_analysis
                }
            ))
        
        for future in as_completed(futures):
            try:
                results.append(future.result())
            except Exception as e:
                st.warning(f"Question processing failed: {str(e)}")
    
    return results

def process_single_question(q: dict, index_path: str, doc_metadata: dict) -> dict:
    """Process a single question with security checks and market context"""
    if not vector_db.is_trusted_path(index_path):
        return {
            "question": q["question"],
            "answer": "⚠️ Security violation detected",
            "category": q.get("category", "General")
        }
    
    try:
        top_chunks = vector_db.search_faiss(q["question"], index_path, k=3)
            
        context = "\n\n".join([
            f"[[Chunk {i+1}]]\n{chunk['text']}\n(Source: {chunk['metadata'].get('source', 'unknown')})"
            for i, chunk in enumerate(top_chunks)
        ])
        
        market_context = ""
        if 'entities' in doc_metadata and 'crypto_project' in doc_metadata['entities']:
            for project in doc_metadata['entities']['crypto_project']:
                project_lower = project.lower()
                for asset, metrics in doc_metadata.get('market_analysis', {}).items():
                    if asset.lower() in project_lower or project_lower in asset.lower():
                        market_context += (
                            f"\n\nMarket Data for {asset}: "
                            f"Current Price: ${metrics.get('current_price', 0):,.2f}, "
                            f"30D Volatility: {metrics.get('30_day_volatility', 0):.2%}, "
                            f"30D Return: {metrics.get('last_30_day_return', 0):.2%}"
                        )
        
        prompt = generate_concise_prompt(
            question=q["question"],
            context=context + market_context,
            metadata={
                **doc_metadata,
                "top_chunk_metadata": top_chunks[0]["metadata"] if top_chunks else {}
            }
        )
        
        answer = llama_chat(prompt)
        evaluation = evaluate_answer(q["question"], answer)
        
        return {
            "question": q["question"],
            "answer": answer,
            "category": q.get("category", "General"),
            "context_used": [chunk["metadata"]["chunk_id"] for chunk in top_chunks] if top_chunks else [],
            "evaluation": evaluation
        }
    except Exception as e:
        return {
            "question": q["question"],
            "answer": f"⚠️ Processing error: {str(e)}",
            "category": q.get("category", "General")
        }

def log_feedback(question: str, answer: str, rating: int, feedback_text: str, filename: str):
    """Log user feedback to CSV for analysis and model improvement"""
    feedback_path = Path("data/feedback/feedback_logs.csv")
    feedback_path.parent.mkdir(exist_ok=True)
    
    feedback_data = {
        "timestamp": datetime.now().isoformat(),
        "question": question,
        "answer": answer,
        "rating": rating,
        "feedback": feedback_text,
        "document": filename,
        "user_agent": str(st.experimental_get_query_params().get("user_agent", ["unknown"])[0])
    }
    
    try:
        if feedback_path.exists():
            existing_df = pd.read_csv(feedback_path)
            updated_df = pd.concat([existing_df, pd.DataFrame([feedback_data])], ignore_index=True)
            updated_df.to_csv(feedback_path, index=False)
        else:
            pd.DataFrame([feedback_data]).to_csv(feedback_path, index=False)
    except Exception as e:
        st.error(f"Failed to log feedback: {str(e)}")

# ---------------------------
# 📌 File Upload & Processing
# ---------------------------
uploaded_file = st.file_uploader("Upload a crypto fund PDF", type="pdf")
if uploaded_file:
    filename = uploaded_file.name
    filepath = os.path.join(UPLOAD_DIR, filename)
    
    with open(filepath, "wb") as f:
        f.write(uploaded_file.getbuffer())
    st.success(f"✅ {filename} uploaded successfully!")

    # ---------------------------
    # 🦠 Text Processing Pipeline
    # ---------------------------
    with st.spinner("🔍 Extracting and cleaning text..."):
        try:
            raw_text = extract_text_from_pdf(filepath)
            if not raw_text:
                st.error("❌ Failed to extract text from PDF")
                st.stop()
                
            cleaned_text = clean_text(raw_text)
            text_path = os.path.join(UPLOAD_DIR, f"{Path(filename).stem}_cleaned.txt")
            with open(text_path, "w", encoding="utf-8") as f:
                f.write(cleaned_text)
                
        except Exception as e:
            st.error(f"❌ Text processing failed: {str(e)}")
            st.stop()

    # ---------------------------
    # ⚠️ Risk Analysis
    # ---------------------------
    with st.spinner("🔍 Analyzing risks..."):
        try:
            risk_terms, risk_score = extract_risk_features(cleaned_text)
            risk_path = os.path.join(UPLOAD_DIR, f"{Path(filename).stem}_risk.json")
            with open(risk_path, "w", encoding="utf-8") as f:
                json.dump({
                    "risk_score": risk_score,
                    "risk_terms": risk_terms,
                    "analysis_time": datetime.now().isoformat()
                }, f, indent=2)
            
            st.subheader("⚠️ Risk Assessment")
            col1, col2 = st.columns(2)
            with col1:
                st.metric("Risk Score", f"{risk_score}/100")
            with col2:
                if risk_score >= 70:
                    st.error("🔴 High Risk")
                elif risk_score >= 40:
                    st.warning("🟠 Moderate Risk")
                else:
                    st.success("🟢 Low Risk")
            
            with st.expander("View Risk Details"):
                for category, terms in risk_terms.items():
                    if terms:
                        st.markdown(f"**{category.title()} Risks**: {', '.join(terms)}")
                        
        except Exception as e:
            st.warning(f"⚠️ Risk analysis incomplete: {str(e)}")
            risk_terms, risk_score = {}, 0

    # ---------------------------
    # 🗃️ Entity Extraction
    # ---------------------------
    with st.spinner("🔍 Identifying key entities..."):
        try:
            named_entities = extract_named_entities(cleaned_text)
            named_entities["crypto_project"] = [
                p for p in named_entities["crypto_project"] 
                if not any(other != p and p in other 
                        for other in named_entities["crypto_project"])
            ]
            
            entities_path = os.path.join(UPLOAD_DIR, f"{Path(filename).stem}_entities.json")
            with open(entities_path, "w", encoding="utf-8") as f:
                json.dump({
                    "entities": named_entities,
                    "extraction_time": datetime.now().isoformat()
                }, f, indent=2)
            
            st.session_state['current_entities'] = named_entities
            
            st.subheader("🧠 Key Entities Identified")
            cols = st.columns(3)
            for i, (ent_type, values) in enumerate(named_entities.items()):
                if values:
                    with cols[i % 3]:
                        with st.expander(f"{ent_type.upper()} ({len(values)})"):
                            st.write(", ".join(values[:15]))
        except Exception as e:
            st.warning(f"⚠️ Entity extraction incomplete: {str(e)}")
            named_entities = {}

    # ---------------------------
    # 📦 Chunking & Indexing
    # ---------------------------
    with st.spinner("📦 Preparing document for analysis..."):
        try:
            chunker = SmartChunker(chunk_size=512, chunk_overlap=64)
            chunks = chunker.chunk_text(
                cleaned_text,
                metadata={
                    "source": filename,
                    "upload_time": datetime.now().isoformat(),
                    "risk_score": risk_score,
                    "entities": named_entities
                }
            )
            
            for i, chunk in enumerate(chunks):
                chunk['metadata']['chunk_id'] = str(uuid.uuid4())
            
            index_path = os.path.join(UPLOAD_DIR, f"{Path(filename).stem}.faiss")
            vector_db.save_to_faiss(chunks, index_path)
            st.success("✅ Document indexed successfully!")
        except Exception as e:
            st.error(f"❌ Failed to index document: {str(e)}")
            st.stop()

    # ---------------------------
    # 🤖 Q&A Processing with Market Context
    # ---------------------------
    st.subheader("💡 Due Diligence Q&A")
    questions = question_bank[:MAX_QUESTIONS]
    
    with st.spinner(f"🤖 Processing {len(questions)} questions..."):
        results = process_question_batch(
            questions,
            index_path,
            {
                "document": filename,
                "risk_score": risk_score,
                "analysis_date": datetime.now().date().isoformat(),
                "entities": named_entities
            }
        )

    # ---------------------------
    # 📊 Results Display
    # ---------------------------
    display_market_insights()
    
    st.subheader("📋 Analysis Results")

    tab1, tab2 = st.tabs(["Q&A Results", "Full Report"])

    with tab1:
        for r in results:
            with st.expander(f"❓ {r['question']} ({r['category']})"):
                st.markdown(f"**Answer:** {r['answer']}")
                
                if "evaluation" in r:
                    st.caption(f"🔍 Evaluation: Accuracy {r['evaluation']['accuracy']}/5 | Completeness {r['evaluation']['completeness']}/5")
                    st.caption(f"Tags: {', '.join(r['evaluation']['tags'])}")
                    st.caption(f"💡 Suggestion: {r['evaluation']['suggestion']}")
                
                with st.form(key=f"feedback_{r['question']}"):
                    rating = st.slider("Rate this answer", 1, 5, 3, key=f"rating_{r['question']}")
                    feedback = st.text_area("Your feedback", key=f"text_{r['question']}")
                    if st.form_submit_button("Submit Feedback"):
                        log_feedback(
                            question=r["question"],
                            answer=r["answer"],
                            rating=rating,
                            feedback_text=feedback,
                            filename=filename
                        )
                        st.success("Feedback saved!")

                if "context_used" in r and r['context_used']:
                    unique_refs = list(set(r['context_used']))[:3]
                    st.caption(f"📌 Source References: {', '.join(unique_refs)}")

    with tab2:
        # Generate PPTX with individual slides per Q&A
        def generate_qa_slides(qa_results, risk_score, risk_terms, entities):
            prs = Presentation()
            
            # Title Slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = "Crypto Fund Due Diligence Report"
            subtitle.text = f"Risk Score: {risk_score}/100 | Generated on {datetime.now().date()}"
            
            # Risk Overview Slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = "Risk Analysis"
            content = slide.placeholders[1]
            
            risk_text = ""
            for category, terms in risk_terms.items():
                if terms:
                    risk_text += f"\n• {category.title()}: {', '.join(terms)}"
            
            content.text = f"Key Risk Factors:{risk_text}"
            
            # Add individual Q&A slides
            for i, qa in enumerate(qa_results):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = f"Q{i+1}: {qa['question'][:50]}..." if len(qa['question']) > 50 else f"Q{i+1}: {qa['question']}"
                
                answer_text = f"Answer: {qa['answer']}\n\n"
                
                if "evaluation" in qa:
                    eval_text = (
                        f"Evaluation:\n"
                        f"- Accuracy: {qa['evaluation']['accuracy']}/5\n"
                        f"- Completeness: {qa['evaluation']['completeness']}/5\n"
                        f"- Suggestions: {qa['evaluation']['suggestion']}\n\n"
                    )
                    answer_text += eval_text
                
                if "context_used" in qa and qa['context_used']:
                    answer_text += "Source References:\n• " + "\n• ".join(qa['context_used'][:3])
                
                content.text = answer_text
            
            # Save to bytes using built-in io
            ppt_stream = io.BytesIO()
            prs.save(ppt_stream)
            ppt_stream.seek(0)
            return ppt_stream.getvalue()

        ppt_bytes = generate_qa_slides(
            qa_results=results,
            risk_score=risk_score,
            risk_terms=risk_terms,
            entities=named_entities
        )
        
        # Save and offer download
        ppt_path = Path("data/processed") / f"due_diligence_{Path(filename).stem}.pptx"
        ppt_path.parent.mkdir(exist_ok=True)
        with open(ppt_path, "wb") as f:
            f.write(ppt_bytes)
        
        st.download_button(
            "📊 Download PPT Report",
            data=ppt_bytes,
            file_name=f"due_diligence_{Path(filename).stem}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    st.success("🎉 Analysis completed successfully!")
else:
    display_market_insights()