import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pymongo import MongoClient
import ollama

# === MongoDB Setup ===
client = MongoClient("mongodb://localhost:27017/")
db = client["crypto_dd"]
funds_collection = db["funds"]

# === LLM Prompting ===
def query_llm(prompt: str, model="llama3.1"):
    response = ollama.chat(model=model, messages=[
        {"role": "system", "content": "You are a professional crypto fund analyst writing due diligence reports for investors."},
        {"role": "user", "content": prompt}
    ])
    return response["message"]["content"].strip()

# === Report Generator ===
def main(fund_name: str, output_dir="output"):
    os.makedirs(output_dir, exist_ok=True)
    fund_doc = funds_collection.find_one({"fund_name": fund_name})
    if not fund_doc:
        raise ValueError(f"❌ Fund '{fund_name}' not found in MongoDB!")

    answers = fund_doc.get("qa_answers", {})
    risk_scores = fund_doc.get("risk_score", {})
    full_text = "\n".join(answers.values())
    today = datetime.today().strftime("%Y-%m-%d")

    summary = query_llm(f"Summarize this fund in bullet points:\n\n{full_text}")
    issues = query_llm(f"List 3–5 red flags or weaknesses from this content:\n\n{full_text}")
    recommendation = query_llm(
        f"You are assessing a crypto fund. Based on these risk scores:\n{risk_scores}\n"
        "Give an investment recommendation: Proceed, Monitor, or Avoid — with a short reason."
    )

    prs = Presentation()

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"{fund_name} – Due Diligence Report"
    slide.placeholders[1].text = f"📅 {today}\n🔎 Platform: DueXpert"

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Findings"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for bullet in summary.split("\n"):
        if bullet.strip():
            p = body.add_paragraph()
            p.text = bullet.strip()
            p.font.size = Pt(14)
            p.font.name = "Arial"

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Risk Score Breakdown"
    table = slide.shapes.add_table(len(risk_scores)+1, 2, Inches(0.5), Inches(1.5), Inches(4.5), Inches(2.5)).table
    table.cell(0, 0).text = "Category"
    table.cell(0, 1).text = "Score"
    for i, (cat, score) in enumerate(risk_scores.items(), start=1):
        table.cell(i, 0).text = cat
        table.cell(i, 1).text = f"{score:.2f}"
    chart_data = CategoryChartData()
    chart_data.categories = list(risk_scores.keys())
    chart_data.add_series("Risk", list(risk_scores.values()))
    slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(5.5), Inches(1.5), Inches(4), Inches(3), chart_data)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Final Recommendation"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    rec_p = body.add_paragraph()
    rec_p.text = recommendation.strip()
    rec_p.font.size = Pt(14)
    rec_p.font.name = "Arial"
    issue_p = body.add_paragraph()
    issue_p.text = "⚠️ Key Issues:"
    issue_p.font.bold = True
    for issue in issues.split("\n"):
        if issue.strip():
            p = body.add_paragraph()
            p.text = f"- {issue.strip()}"
            p.font.size = Pt(12)
            p.font.name = "Arial"

    pptx_path = os.path.join(output_dir, f"{fund_name}_DueDiligenceReport.pptx")
    prs.save(pptx_path)
    print(f"✅ Report saved: {pptx_path}")
    return pptx_path