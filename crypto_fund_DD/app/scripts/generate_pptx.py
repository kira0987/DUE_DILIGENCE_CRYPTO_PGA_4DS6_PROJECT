
import os
import json
import pandas as pd
import re
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from collections import defaultdict
import spacy
import numpy as np
import logging
from tqdm import tqdm

# --- Logging Setup ---
os.makedirs("output", exist_ok=True)
logging.basicConfig(
	filename="output/due_diligence.log",
	level=logging.INFO,
	format="%(asctime)s - %(levelname)s - %(message)s",
	filemode="w"
)
logger = logging.getLogger(__name__)

# --- Paths ---
ANSWERS_CSV = "data/auto_answered_questions.csv"
CLASSIFIED_JSON = "data/classified_questions.json"
LOGO_PATH = "assets/dueexpert.png"
OUTPUT_DIR = "output"
COMPANY_LOGO_PATH = "assets/VALUE.png"

# --- Global Variables ---
CRYPTO_COLORS = {
	"primary": RGBColor(255, 255, 255),
	"accent": RGBColor(153, 51, 255),
	"highlight": RGBColor(54, 95, 199),
	"text": RGBColor(0, 0, 0),
	"background": RGBColor(255, 255, 255),
}

TAG_OBJECTIVES = {
	"AML / KYC": "Ensure robust Anti-Money Laundering and Know Your Customer compliance.",
	"Community & UX": "Assess user engagement, usability, and community transparency.",
	"Custody & Asset Security": "Evaluate asset storage and custody security mechanisms.",
	"Cybersecurity & Data Privacy": "Examine cybersecurity and data protection practices.",
	"ESG & Sustainability": "Evaluate environmental, social, and governance sustainability.",
	"Financial Health": "Analyze financial statements and health indicators.",
	"Governance": "Assess decision-making and board oversight structures.",
	"IP & Contracts": "Evaluate intellectual property and contract integrity.",
	"Legal & Regulatory": "Ensure legal and regulatory compliance.",
	"Risk Management": "Identify risk management strategies.",
	"Strategy & Competitive Positioning": "Understand strategy and market positioning.",
	"Technology & Infrastructure": "Evaluate technological scalability and reliability.",
	"Tokenomics & Trading Integrity": "Assess token economics and trading transparency.",
	"Future Outlook": "Evaluate alignment with 2025 crypto trends and innovations."
}

DUE_DILIGENCE_CRITERIA = {
	"AML / KYC": ["AML policies", "KYC procedures", "compliance measures", "regulatory adherence"],
	"Community & UX": ["user engagement", "platform usability", "community transparency", "feedback mechanisms"],
	"Custody & Asset Security": ["asset storage", "multi-signature wallets", "asset ownership verification", "insurance coverage"],
	"Cybersecurity & Data Privacy": ["blockchain technology", "cybersecurity measures", "two-factor authentication", "data protection policy"],
	"ESG & Sustainability": ["sustainability practices", "energy consumption", "environmental impact", "sustainability audits"],
	"Financial Health": ["financial statements", "cash flow projections", "tax disputes", "annual auditing"],
	"Governance": ["board of directors", "conflict of interest policies", "external audits", "succession planning"],
	"IP & Contracts": ["intellectual property assets", "IP infringement", "smart contract audits", "major contracts"],
	"Legal & Regulatory": ["licenses and permits", "international law compliance", "sanctions monitoring", "regulatory filings"],
	"Risk Management": ["risk management strategies", "market volatility", "business continuity plan", "operational resilience"],
	"Strategy & Competitive Positioning": ["market positioning", "competitor analysis", "investment strategies", "growth strategies"],
	"Technology & Infrastructure": ["scalability plans", "security measures", "blockchain technology", "infrastructure upgrades"],
	"Tokenomics & Trading Integrity": ["securities compliance", "trade surveillance", "valuation consistency", "code reviews"],
	"Future Outlook": ["stablecoin compliance", "AI integration", "institutional adoption", "quantum resistance"]
}

IMPORTANCE_WEIGHTS = {
	"AML / KYC": 1.0,
	"Legal & Regulatory": 1.0,
	"Cybersecurity & Data Privacy": 0.8,
	"Custody & Asset Security": 0.8,
	"Financial Health": 0.64,
	"Governance": 0.48,
	"Risk Management": 0.48,
	"Tokenomics & Trading Integrity": 0.36,
	"Community & UX": 0.24,
	"ESG & Sustainability": 0.16,
	"IP & Contracts": 0.16,
	"Strategy & Competitive Positioning": 0.16,
	"Technology & Infrastructure": 0.24,
	"Future Outlook": 0.08
}

nlp = spacy.load("en_core_web_sm")

# --- Helper Functions ---
def get_next_output_filename(base_name):
	os.makedirs(OUTPUT_DIR, exist_ok=True)
	i = 0
	while True:
		suffix = f" ({i})" if i else ""
		path = os.path.join(OUTPUT_DIR, f"{base_name}{suffix}.pptx")
		if not os.path.exists(path):
			return path
		i += 1

def load_data():
	try:
		df = pd.read_csv(ANSWERS_CSV)
		df.columns = [c.lower() for c in df.columns]
		df["question"] = df["question"].astype(str).str.strip()
		df["summary"] = df["answer"].apply(extract_summary)
		metric_columns = ["relevance", "completeness", "clarity", "faithfulness"]
		for col in metric_columns:
			if col in df.columns:
				df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)
				df[col] = df[col].replace([np.inf, -np.inf], 0)
		with open(CLASSIFIED_JSON, "r", encoding="utf-8") as f:
			classified = json.load(f)
		logger.info(f"Data loaded: {len(df)} questions, {len(classified)} classified entries")
		return df, classified
	except FileNotFoundError as e:
		logger.error(f"CSV or JSON file not found: {str(e)}")
		print("Warning: CSV or JSON file not found. Using fallback data.")
		df = pd.DataFrame({
			"question": ["What are the AML policies?", "Who are the managing partners?"],
			"answer": ["The fund has robust AML policies.", "Managing partners oversee operations."],
			"summary": ["The fund has robust AML policies.", "Managing partners oversee operations."],
			"relevance": [0.9, 0.8],
			"completeness": [0.85, 0.75],
			"clarity": [0.9, 0.85],
			"faithfulness": [0.95, 0.9]
		})
		classified = [
			{"question": "What are the AML policies?", "tag": "AML / KYC"},
			{"question": "Who are the managing partners?", "tag": "Governance"}
		]
		logger.info("Fallback data loaded")
		return df, classified

def extract_fund_name(df):
	pattern = r'\b[\w\s]+?\sFund\b'
	for text in pd.concat([df["question"].astype(str), df["answer"].astype(str)], ignore_index=True):
		if pd.isna(text):
			continue
		match = re.search(pattern, text)
		if match:
			fund_name = match.group(0).rstrip(" Fund").strip()
			logger.info(f"Fund name extracted: {fund_name}")
			return fund_name
	logger.warning("No fund name found, using default")
	return "Crypto Fund"

def extract_summary(text):
	if not isinstance(text, str) or len(text.strip()) < 20:
		logger.warning(f"Invalid summary text: {text[:30]}...")
		return None
	lower = text.lower()
	if any(x in lower for x in ["not mentioned", "cannot determine", "unavailable"]):
		logger.warning(f"Negative summary detected: {text[:30]}...")
		return None
	clean = text.replace("Based on the document,", "").replace("According to", "").strip()
	summary = clean.split(". ")[0].strip().capitalize() + "."
	logger.info(f"Summary extracted: {summary[:30]}...")
	return summary

def apply_crypto_theme(slide):
	bg = slide.background
	fill = bg.fill
	fill.gradient()
	fill.gradient_stops[0].color.rgb = CRYPTO_COLORS["background"]
	fill.gradient_stops[1].color.rgb = CRYPTO_COLORS["primary"]
	logger.info("Crypto theme applied to slide")

def add_branding(slide, current_tag="General"):
	if os.path.exists(LOGO_PATH):
		slide.shapes.add_picture(LOGO_PATH, Inches(0.3), Inches(0.1), height=Inches(0.6))
	if os.path.exists(COMPANY_LOGO_PATH):
		slide.shapes.add_picture(COMPANY_LOGO_PATH, Inches(10.3), Inches(6.7), height=Inches(0.6))
	footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9), Inches(0.3)).text_frame
	footer.text = f"DueXpert – AI Crypto Fund Due Diligence Suite | {current_tag}"
	footer.paragraphs[0].font.size = Pt(10)
	footer.paragraphs[0].font.color.rgb = CRYPTO_COLORS["text"]
	footer.paragraphs[0].font.name = "Montserrat"
	footer.paragraphs[0].font.bold = True
	logger.info(f"Branding added for tag: {current_tag}")

# --- Analysis Functions ---
def detect_negative_sentiment(text):
	doc = nlp(text.lower())
	negative_indicators = ["not mentioned", "lacking", "insufficient", "unavailable", "no details", "unknown", "weak", "incomplete"]
	result = any(indicator in text.lower() for indicator in negative_indicators)
	logger.info(f"Negative sentiment check for '{text[:30]}...': {result}")
	return result

def analyze_tag_text(tag, findings, issues):
	tag_text = " ".join(findings) + " " + " ".join(issues)
	criteria = DUE_DILIGENCE_CRITERIA[tag]
	doc = nlp(tag_text.lower())
	found_criteria = []
	for criterion in criteria:
		criterion_doc = nlp(criterion.lower())
		similarity = doc.similarity(criterion_doc)
		keyword_match = criterion.lower() in tag_text.lower()
		if similarity >= 0.7 or keyword_match:
			found_criteria.append(criterion)
	found_criteria = list(set(found_criteria))
	completeness = len(found_criteria) / len(criteria) * 100 if criteria else 0
	missing = [c for c in criteria if c not in found_criteria]
	if detect_negative_sentiment(tag_text):
		completeness *= 0.8
	logger.info(f"Tag {tag} analysis: completeness={completeness:.1f}%, found={len(found_criteria)}, missing={len(missing)}")
	return {
		"completeness": completeness,
		"found": found_criteria,
		"missing": missing
	}

def calculate_question_risk(row, tag, tag_weight):
	"""Calculate risk score for a single question based on metrics."""
	metrics = ["relevance", "completeness", "clarity", "faithfulness"]
	metric_scores = [float(row[m]) if pd.notna(row[m]) and np.isfinite(row[m]) else 0.0 for m in metrics]
	avg_metric = np.mean(metric_scores)
	
	answer_score = 0.5  # Fallback score
	
	combined_score = (avg_metric + answer_score) / 2
	risk = (1 - combined_score) * tag_weight * 100
	if detect_negative_sentiment(row["answer"]):
		risk *= 1.2
	
	final_risk = float(min(risk, 100))
	logger.info(f"Question risk for '{row['question'][:30]}...': {final_risk:.1f}, Answer Score: {answer_score:.2f}")
	return final_risk

def calculate_risk_score(analysis, df, classified):
	"""Calculate tag-level risk scores, incorporating question-level risks."""
	risk_scores = {}
	tag_map = {q["question"]: q["tag"] for q in classified}
	df["tag"] = df["question"].map(tag_map)
	
	for tag in analysis:
		tag_df = df[df["tag"] == tag].copy()
		weight = IMPORTANCE_WEIGHTS[tag]
		completeness = analysis[tag]["completeness"] / 100
		
		question_risks = []
		for _, row in tag_df.iterrows():
			question_risk = calculate_question_risk(row, tag, weight)
			question_risks.append(question_risk)
		
		tag_risk = (1 - completeness) * weight * 100
		if question_risks:
			avg_question_risk = np.mean(question_risks)
			tag_risk = (tag_risk + avg_question_risk) / 2
		if len(analysis[tag]["missing"]) > len(analysis[tag]["found"]):
			tag_risk *= 1.2
		risk_scores[tag] = min(tag_risk, 100)
		logger.info(f"Tag {tag} risk score: {risk_scores[tag]:.1f}")
	
	return risk_scores

# --- Analytics PPTX Functions ---
def create_analytics_cover_slide(prs, fund_name):
	slide = prs.slides.add_slide(prs.slide_layouts[0])
	apply_crypto_theme(slide)
	title = slide.shapes.title
	title.text = f"{fund_name} Due Diligence Analytics"
	title.text_frame.paragraphs[0].font.size = Pt(36)
	title.text_frame.paragraphs[0].font.bold = True
	title.text_frame.paragraphs[0].font.color.rgb = CRYPTO_COLORS["highlight"]
	title.text_frame.paragraphs[0].font.name = "Montserrat"
	subtitle = slide.placeholders[1]
	subtitle.text = f"Prepared by DueXpert AI | {datetime.today().strftime('%Y-%m-%d')}"
	subtitle.text_frame.paragraphs[0].font.size = Pt(20)
	subtitle.text_frame.paragraphs[0].font.color.rgb = CRYPTO_COLORS["text"]
	subtitle.text_frame.paragraphs[0].font.name = "Montserrat"
	add_branding(slide)
	logger.info("Cover slide created")
	return slide

def create_analytics_toc_slide(prs, tag_page_numbers):
	slide = prs.slides.add_slide(prs.slide_layouts[6])
	apply_crypto_theme(slide)
	add_branding(slide, current_tag="Table of Contents")
	title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.5)).text_frame
	title.text = "Table of Contents"
	title.paragraphs[0].font.size = Pt(28)
	title.paragraphs[0].font.bold = True
	title.paragraphs[0].font.color.rgb = CRYPTO_COLORS["highlight"]
	title.paragraphs[0].font.name = "Montserrat"
	
	content = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.5), Inches(5.5)).text_frame
	content.word_wrap = True
	p = content.add_paragraph()
	p.text = "2. Key Summary and Findings"
	p.font.size = Pt(16)
	p.font.color.rgb = CRYPTO_COLORS["text"]
	p.font.name = "Montserrat"
	p = content.add_paragraph()
	p.text = "3. Issues Faced"
	p.font.size = Pt(16)
	p.font.color.rgb = CRYPTO_COLORS["text"]
	p.font.name = "Montserrat"
	p = content.add_paragraph()
	p.text = "4. Risk Scoring"
	p.font.size = Pt(16)
	p.font.color.rgb = CRYPTO_COLORS["text"]
	p.font.name = "Montserrat"
	p = content.add_paragraph()
	p.text = "5. Risk Visualization"
	p.font.size = Pt(16)
	p.font.color.rgb = CRYPTO_COLORS["text"]
	p.font.name = "Montserrat"
	p = content.add_paragraph()
	p.text = "6. Recommendations"
	p.font.size = Pt(16)
	p.font.color.rgb = CRYPTO_COLORS["text"]
	p.font.name = "Montserrat"
	logger.info("TOC slide created")
	return slide

def create_key_findings_summary_slide(prs, df, classified, page_num):
	slide = prs.slides.add_slide(prs.slide_layouts[6])
	apply_crypto_theme(slide)
	add_branding(slide, current_tag="Key Findings Summary")
	
	title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.5)).text_frame
	title.text = f"{page_num}. Key Summary and Findings"
	title.paragraphs[0].font.size = Pt(24)
	title.paragraphs[0].font.bold = True
	title.paragraphs[0].font.color.rgb = CRYPTO_COLORS["highlight"]
	title.paragraphs[0].font.name = "Montserrat"
	
	tag_map = {q["question"]: q["tag"] for q in classified}
	df["tag"] = df["question"].map(tag_map)
	tag_groups = df.groupby("tag").agg({
		"summary": lambda x: [s for s in x if s],
		"relevance": "mean",
		"completeness": "mean",
		"clarity": "mean",
		"faithfulness": "mean"
	}).reset_index()
	
	findings = []
	for _, row in tag_groups.iterrows():
		tag = row["tag"]
		if row["summary"]:
			findings.append(f"{tag}: {'; '.join(row['summary'][:3])}{'...' if len(row['summary']) > 3 else ''}")
	
	content_text = "No findings available." if not findings else "\n".join(findings[:10])
	content = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(10.0), Inches(5.5)).text_frame
	content.text = content_text
	content.word_wrap = True
	for p in content.paragraphs:
		p.font.size = Pt(12)
		p.font.color.rgb = CRYPTO_COLORS["text"]
		p.font.name = "Montserrat"
	
	page = slide.shapes.add_textbox(Inches(9.0), Inches(7.0), Inches(1), Inches(0.3)).text_frame
	page.text = f"Page {page_num}"
	page.paragraphs[0].font.size = Pt(10)
	page.paragraphs[0].font.color.rgb = CRYPTO_COLORS["text"]
	page.paragraphs[0].font.name = "Montserrat"
	page.paragraphs[0].font.bold = True
	logger.info("Key Findings Summary slide created")
	return slide

def create_issues_faced_slide(prs, tag_findings, df, page_num):
	slide = prs.slides.add_slide(prs.slide_layouts[6])
	apply_crypto_theme(slide)
	add_branding(slide, current_tag="Issues Faced")
	
	title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.5)).text_frame
	title.text = f"{page_num}. Issues Faced"
	title.paragraphs[0].font.size = Pt(24)
	title.paragraphs[0].font.bold = True
	title.paragraphs[0].font.color.rgb = CRYPTO_COLORS["highlight"]
	title.paragraphs[0].font.name = "Montserrat"
	
	issues_text = []
	for tag, (findings, issues) in tag_findings.items():
		if issues:
			issues_text.append(f"{tag}: {', '.join(issues[:3])}{'...' if len(issues) > 3 else ''}")
		for _, row in df[df["tag"] == tag].iterrows():
			if detect_negative_sentiment(row["answer"]):
				issues_text.append(f"{tag}: Negative response for '{row['question'][:50]}...'")
	
	content_text = "No issues detected." if not issues_text else "\n".join(issues_text[:10])
	content = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(10.0), Inches(5.5)).text_frame
	content.text = content_text
	content.word_wrap = True
	for p in content.paragraphs:
		p.font.size = Pt(12)
		p.font.color.rgb = CRYPTO_COLORS["text"]
		p.font.name = "Montserrat"
	
	page = slide.shapes.add_textbox(Inches(9.0), Inches(7.0), Inches(1), Inches(0.3)).text_frame
	page.text = f"Page {page_num}"
	page.paragraphs[0].font.size = Pt(10)
	page.paragraphs[0].font.color.rgb = CRYPTO_COLORS["text"]
	page.paragraphs[0].font.name = "Montserrat"
	page.paragraphs[0].font.bold = True
	
	logger.info("Issues Faced slide created")
	return slide

def create_risk_scoring_slide(prs, risk_scores, page_num):
	slide = prs.slides.add_slide(prs.slide_layouts[6])
	apply_crypto_theme(slide)
	add_branding(slide, current_tag="Risk Scoring")
	
	title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.5)).text_frame
	title.text = f"{page_num}. Risk Scoring"
	title.paragraphs[0].font.size = Pt(24)
	title.paragraphs[0].font.bold = True
	title.paragraphs[0].font.color.rgb = CRYPTO_COLORS["highlight"]
	title.paragraphs[0].font.name = "Montserrat"
	
	high_risk = sum(1 for tag in risk_scores if risk_scores[tag] >= 60)
	medium_risk = sum(1 for tag in risk_scores if 30 <= risk_scores[tag] < 60)
	low_risk = sum(1 for tag in risk_scores if risk_scores[tag] < 30)
	avg_risk = sum(risk_scores[tag] for tag in risk_scores) / len(risk_scores)
	critical_gaps = [tag for tag in risk_scores if IMPORTANCE_WEIGHTS[tag] >= 0.8 and risk_scores[tag] >= 60]
	
	content_text = (
		f"Risk Scores by Category:\n"
		+ "\n".join(f"- {tag}: {risk_scores[tag]:.1f}" for tag in sorted(risk_scores))
		+ f"\n\nSummary:\n"
		f"- Total Topics: {len(risk_scores)}\n"
		f"- High Risk (>=60): {high_risk}\n"
		f"- Medium Risk (30-60): {medium_risk}\n"
		f"- Low Risk (<30): {low_risk}\n"
		f"- Avg Risk Score: {avg_risk:.1f}\n"
		f"- Critical Gaps: {', '.join(critical_gaps) if critical_gaps else 'None'}"
	)
	content = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(10.0), Inches(5.5)).text_frame
	content.text = content_text
	content.word_wrap = True
	for p in content.paragraphs:
		p.font.size = Pt(12)
		p.font.color.rgb = CRYPTO_COLORS["text"]
		p.font.name = "Montserrat"
	
	page = slide.shapes.add_textbox(Inches(9.0), Inches(7.0), Inches(1), Inches(0.3)).text_frame
	page.text = f"Page {page_num}"
	page.paragraphs[0].font.size = Pt(10)
	page.paragraphs[0].font.color.rgb = CRYPTO_COLORS["text"]
	page.paragraphs[0].font.name = "Montserrat"
	page.paragraphs[0].font.bold = True
	
	logger.info("Risk Scoring slide created")
	return slide

def create_risk_visualization_slide(prs, analysis, risk_scores, page_num):
	slide = prs.slides.add_slide(prs.slide_layouts[6])
	apply_crypto_theme(slide)
	add_branding(slide, current_tag="Risk Visualization")
	
	title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.5)).text_frame
	title.text = f"{page_num}. Risk Visualization"
	title.paragraphs[0].font.size = Pt(24)
	title.paragraphs[0].font.bold = True
	title.paragraphs[0].font.color.rgb = CRYPTO_COLORS["highlight"]
	title.paragraphs[0].font.name = "Montserrat"
	
	try:
		chart_data = CategoryChartData()
		chart_data.categories = sorted(analysis.keys())
		chart_data.add_series("Risk Score", [risk_scores[tag] for tag in sorted(analysis)])
		x, y, cx, cy = Inches(0.8), Inches(1.2), Inches(10.0), Inches(5.0)
		chart = slide.shapes.add_chart(XL_CHART_TYPE.AREA, x, y, cx, cy, chart_data).chart
		chart.has_legend = True
		chart.legend.position = XL_LEGEND_POSITION.BOTTOM
		chart.chart_title.text_frame.text = "Risk Scores by Topic (0-100)"
		chart.category_axis.tick_labels.font.size = Pt(10)
		logger.info("Risk Visualization chart created")
	except Exception as e:
		logger.error(f"Failed to create Risk Visualization chart: {str(e)}")
		content = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(10.0), Inches(5.5)).text_frame
		content.text = "Error: Unable to generate risk visualization chart."
		content.word_wrap = True
		for p in content.paragraphs:
			p.font.size = Pt(12)
			p.font.color.rgb = CRYPTO_COLORS["text"]
			p.font.name = "Montserrat"
	
	page = slide.shapes.add_textbox(Inches(9.0), Inches(7.0), Inches(1), Inches(0.3)).text_frame
	page.text = f"Page {page_num}"
	page.paragraphs[0].font.size = Pt(10)
	page.paragraphs[0].font.color.rgb = CRYPTO_COLORS["text"]
	page.paragraphs[0].font.name = "Montserrat"
	page.paragraphs[0].font.bold = True
	
	logger.info("Risk Visualization slide created")
	return slide

def create_recommendations_slide(prs, analysis, risk_scores, page_num):
	slide = prs.slides.add_slide(prs.slide_layouts[6])
	apply_crypto_theme(slide)
	add_branding(slide, current_tag="Recommendations")
	
	title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.5)).text_frame
	title.text = f"{page_num}. Recommendations"
	title.paragraphs[0].font.size = Pt(24)
	title.paragraphs[0].font.bold = True
	title.paragraphs[0].font.color.rgb = CRYPTO_COLORS["highlight"]
	title.paragraphs[0].font.name = "Montserrat"
	
	recommendations = []
	for tag in sorted(analysis.keys()):
		if risk_scores[tag] >= 60 and IMPORTANCE_WEIGHTS[tag] >= 0.8:
			recommendations.append(f"- Strengthen {tag}: High risk detected ({risk_scores[tag]:.1f}). Address missing criteria: {', '.join(analysis[tag]['missing'][:3])}{'...' if len(analysis[tag]['missing']) > 3 else ''}.")
		elif risk_scores[tag] >= 30:
			recommendations.append(f"- Review {tag}: Moderate risk ({risk_scores[tag]:.1f}). Ensure compliance with {', '.join(analysis[tag]['missing'][:2])}{'...' if len(analysis[tag]['missing']) > 2 else ''}.")
	
	content_text = "No recommendations needed." if not recommendations else "\n".join(recommendations[:10])
	content = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(10.0), Inches(5.5)).text_frame
	content.text = content_text
	content.word_wrap = True
	for p in content.paragraphs:
		p.font.size = Pt(12)
		p.font.color.rgb = CRYPTO_COLORS["text"]
		p.font.name = "Montserrat"
	
	page = slide.shapes.add_textbox(Inches(9.0), Inches(7.0), Inches(1), Inches(0.3)).text_frame
	page.text = f"Page {page_num}"
	page.paragraphs[0].font.size = Pt(10)
	page.paragraphs[0].font.color.rgb = CRYPTO_COLORS["text"]
	page.paragraphs[0].font.name = "Montserrat"
	page.paragraphs[0].font.bold = True
	
	logger.info("Recommendations slide created")
	return slide

# --- Main Function ---
def main():
	try:
		logger.info("Starting report generation")
		df, classified = load_data()
		
		tag_map = {q["question"]: q["tag"] for q in classified}
		df["tag"] = df["question"].map(tag_map)
		logger.info("Calculating answer quality scores")
		df["answer_quality"] = 0.5  # Fallback value for answer quality
		
		tag_map = defaultdict(list)
		for q in classified:
			if isinstance(q, dict) and "tag" in q and "question" in q:
				tag_map[q["tag"]].append(q)
			else:
				logger.warning(f"Invalid question entry skipped: {q}")
		
		tag_summary = {}
		tag_findings = {}
		for tag in sorted(tag_map):
			findings, issues = [], []
			for q in tag_map[tag]:
				qt = q["question"].strip()
				row = df[df["question"] == qt]
				if not row.empty:
					summary = extract_summary(row.iloc[0]["answer"])
					if summary and not detect_negative_sentiment(row.iloc[0]["answer"]):
						findings.append(summary)
					else:
						issues.append(qt)
				else:
					issues.append(qt)
			tag_summary[tag] = (len(findings), len(issues))
			tag_findings[tag] = (findings, issues)
		
		analysis = {}
		for tag in TAG_OBJECTIVES:
			if tag in tag_findings:
				findings, issues = tag_findings[tag]
				analysis[tag] = analyze_tag_text(tag, findings, issues)
			else:
				analysis[tag] = {"completeness": 0, "found": [], "missing": DUE_DILIGENCE_CRITERIA[tag]}
		
		risk_scores = calculate_risk_score(analysis, df, classified)
		analytics_path = get_next_output_filename("due_diligence_report")
		prs_analytics = Presentation()
		prs_analytics.slide_width = Inches(12)
		prs_analytics.slide_height = Inches(7.5)
		
		fund_name = extract_fund_name(df)
		cover_slide = create_analytics_cover_slide(prs_analytics, fund_name)
		
		analytics_tag_page_numbers = {
			"Key Summary and Findings": 2,
			"Issues Faced": 3,
			"Risk Scoring": 4,
			"Risk Visualization": 5,
			"Recommendations": 6
		}
		
		toc_slide = create_analytics_toc_slide(prs_analytics, analytics_tag_page_numbers)
		page_counter = 2
		analytics_slides = [cover_slide, toc_slide]
		
		key_findings_slide = create_key_findings_summary_slide(prs_analytics, df, classified, page_counter)
		analytics_slides.append(key_findings_slide)
		page_counter += 1
		
		issues_faced_slide = create_issues_faced_slide(prs_analytics, tag_findings, df, page_counter)
		analytics_slides.append(issues_faced_slide)
		page_counter += 1
		
		risk_scoring_slide = create_risk_scoring_slide(prs_analytics, risk_scores, page_counter)
		analytics_slides.append(risk_scoring_slide)
		page_counter += 1
		
		risk_visualization_slide = create_risk_visualization_slide(prs_analytics, analysis, risk_scores, page_counter)
		analytics_slides.append(risk_visualization_slide)
		page_counter += 1
		
		recommendations_slide = create_recommendations_slide(prs_analytics, analysis, risk_scores, page_counter)
		analytics_slides.append(recommendations_slide)
		
		prs_analytics.save(analytics_path)
		logger.info(f"Analytics Report Generated: {analytics_path}")
		print(f"✅ Analytics Report Generated: {analytics_path}")
	
	except Exception as e:
		logger.error(f"Failed to generate analytics report: {str(e)}")
		print(f"❌ Failed to generate analytics report: {str(e)}")
		raise

if __name__ == "__main__":
	main()
