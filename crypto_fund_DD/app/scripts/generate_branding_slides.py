# generate_branding_slides.py – Embed Header & Footer in Main PPTX (DueXpert Style)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import os

# --- Constants ---
LOGO_PATH = "assets/shield.png"
CONTACT_LINES = [
    "DueXpert – AI Crypto Fund Due Diligence Suite",
    "123 Rue de l’Innovation, 75000 Paris | France",
    "T: +33 1 42 00 00 00 | E: info@duexpert.ai | W: www.duexpert.ai"
]

# --- Embed Header/Footer Utility ---
def add_branding(slide):
    # Header Logo
    slide.shapes.add_picture(LOGO_PATH, Inches(0.3), Inches(0.3), height=Inches(0.7))

    # Footer Text Box
    footer_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.8), Inches(8.5), Inches(1.0))
    tf = footer_box.text_frame
    for line in CONTACT_LINES:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.name = 'Segoe UI'
        p.font.color.rgb = RGBColor(100, 100, 100)
    tf.paragraphs[0].space_after = Pt(6)

# --- Use in Slide Generator ---
# In your main slide generator (e.g., generate_pptx.py), call:
# add_branding(slide)
# at the end of each slide creation block.
