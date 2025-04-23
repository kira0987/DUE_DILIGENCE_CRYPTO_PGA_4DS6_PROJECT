from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import json
import os

def generate_ppt(answers_path: str, output_path: str = "answers/report.pptx"):
    if not os.path.exists(answers_path):
        raise FileNotFoundError(f"❌ Le fichier de réponses est introuvable : {answers_path}")

    with open(answers_path, "r", encoding="utf-8") as f:
        answers = json.load(f)

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]

    # Slide d’intro
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Due Diligence Report"
    slide.placeholders[1].text = "Généré automatiquement à partir du livre blanc"

    # Slides par question
    for answer in answers.values():
        slide_layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = answer["question"]
        content.text = answer["answer"]

    prs.save(output_path)
    print(f"✅ PowerPoint généré avec succès : {output_path}")
